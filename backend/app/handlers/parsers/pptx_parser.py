""".pptx：用 python-pptx 逐页提取标题、正文文本框、表格，保留幻灯片序号。

幻灯片序号保留在文本里（``## 第 N 页``），因为这批语料的 pptx 是操作步骤类
教程（"一、下载并安装APP" "二、打开APP注册账户"……），页码本身就是步骤顺序
的一部分，丢了序号会让分块后的文本失去"这是第几步"的上下文。
"""
from __future__ import annotations

from pptx import Presentation

from ._source import Source, as_path_or_stream
from .markdown_table import rows_to_markdown
from .types import ParseResult


def parse_pptx(source: Source) -> ParseResult:
    try:
        prs = Presentation(as_path_or_stream(source))
    except Exception as e:
        return ParseResult.failure(f"{type(e).__name__}: {e}")

    slide_texts = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = [f"## 第 {i} 页"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
            if shape.has_table:
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                md = rows_to_markdown(rows)
                if md:
                    parts.append(md)
        if len(parts) > 1:  # 除了页码标题还有实际内容
            slide_texts.append("\n\n".join(parts))

    return ParseResult.success("\n\n".join(slide_texts))
