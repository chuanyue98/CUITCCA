"""``handlers/parsers`` 注册表与各格式解析器的单元测试。

这次改造要解决的核心问题是"257 个语料文件里有 54 个（21%）因为格式白名单
被静默跳过"，所以测试重点不只是"能不能解析"，更是**三种解析结局有没有被
如实区分**：成功 / 明确失败 / 能力不可用。静默跳过是被明确禁止的行为，
这里用测试把它钉住。

旧格式（.doc/.xls）没法用纯 Python 现造（xlrd 只读不写，olefile 同理），
所以这两种直接对仓库里的真实语料文件做集成断言——它们是 git 跟踪的，CI 上
同样存在。
"""
import unittest
from io import BytesIO
from pathlib import Path

from docx import Document as DocxDocument
from handlers.parsers import ParseStatus, parse_bytes, parse_path, supported_extensions
from handlers.parsers.markdown_table import rows_to_markdown
from pptx import Presentation

import tests._pathsetup  # noqa: F401

REPO_ROOT = Path(__file__).resolve().parent.parent
CORPUS = REPO_ROOT / "信息搜集汇总"

REAL_DOC = CORPUS / "学生服务" / "毕业证、学位证遗失补办申请表.doc"
REAL_XLS = CORPUS / "就业信息" / "成都信息工程大学2023届学院就业负责人联系方式.xls"
REAL_JPG = CORPUS / "学生服务" / "复学流程.jpg"


class RegistryDispatchTest(unittest.TestCase):
    def test_newly_supported_extensions_are_registered(self):
        """这 5 类就是改造前被白名单挡在门外的格式。"""
        exts = supported_extensions()
        for ext in (".doc", ".xls", ".pptx", ".html", ".htm", ".jpg", ".jpeg", ".png"):
            self.assertIn(ext, exts, f"{ext} 应当已注册")

    def test_unknown_extension_returns_failure_not_exception(self):
        """分派层的契约：永远返回 ParseResult，不往上抛异常。

        批量摄取时一个不认识的文件不能中断整批处理。
        """
        result = parse_bytes(".zip", b"PK\x03\x04")
        self.assertIs(result.status, ParseStatus.FAILURE)
        self.assertIn("zip", result.reason)

    def test_extension_matching_is_case_insensitive(self):
        self.assertTrue(parse_bytes(".TXT", "大写扩展名".encode()).ok)

    def test_parser_internal_exception_is_converted_to_failure(self):
        """解析器实现里没预料到的异常也要被兜住，否则一个坏文件会炸掉整批摄取。"""
        result = parse_bytes(".pdf", b"definitely not a pdf")
        self.assertIs(result.status, ParseStatus.FAILURE)
        self.assertTrue(result.reason)


class TextLikeParserTest(unittest.TestCase):
    def test_txt_falls_back_to_gbk(self):
        result = parse_bytes(".txt", "成都信息工程大学".encode("gbk"))
        self.assertTrue(result.ok)
        self.assertEqual(result.text, "成都信息工程大学")

    def test_txt_reports_failure_on_undecodable_bytes(self):
        result = parse_bytes(".txt", b"\xff\xfe\x00\x01\x02\xff\xfe")
        self.assertIs(result.status, ParseStatus.FAILURE)

    def test_csv_preserves_column_structure(self):
        csv_bytes = "学院,联系人,电话\n大气科学学院,吴老师,18683376821\n".encode()
        result = parse_bytes(".csv", csv_bytes)
        self.assertTrue(result.ok)
        self.assertIn("学院 | 联系人 | 电话", result.text)

    def test_html_strips_noise_tags(self):
        html = b"""<html><head><style>.a{color:red}</style></head>
        <body><nav>\xe5\xaf\xbc\xe8\x88\xaa</nav>
        <p>\xe5\x9b\xbe\xe4\xb9\xa6\xe9\xa6\x86\xe5\xbc\x80\xe6\x94\xbe\xe6\x97\xb6\xe9\x97\xb4 8:00-22:00</p>
        <script>var x=1;</script></body></html>"""
        result = parse_bytes(".html", html)
        self.assertTrue(result.ok)
        self.assertIn("图书馆开放时间 8:00-22:00", result.text)
        self.assertNotIn("var x", result.text)
        self.assertNotIn("color:red", result.text)
        self.assertNotIn("导航", result.text)


class MarkdownTableTest(unittest.TestCase):
    """表格转 Markdown 是 pdf/docx/pptx 三个解析器共用的表达方式——语料里大量
    内容（借阅规则、时刻表、联系方式表）本质是表格，塌成一行就丢了行列对应
    关系，这是真实的检索质量损失。"""

    def test_renders_header_and_rows(self):
        md = rows_to_markdown([["读者类型", "册数"], ["本科生", "10册"]])
        self.assertIn("| 读者类型 | 册数 |", md)
        self.assertIn("| --- | --- |", md)
        self.assertIn("| 本科生 | 10册 |", md)

    def test_empty_input_returns_empty_string(self):
        """空表格不能渲染成"只有表头分隔线"的空壳，调用方要能直接跳过。"""
        self.assertEqual(rows_to_markdown([]), "")
        self.assertEqual(rows_to_markdown([[None, ""], ["", None]]), "")

    def test_cell_content_that_would_break_table_syntax_is_neutralised(self):
        md = rows_to_markdown([["a\nb", "c|d"]])
        self.assertNotIn("\nb", md)
        self.assertIn("c/d", md)

    def test_ragged_rows_are_padded(self):
        md = rows_to_markdown([["a", "b", "c"], ["x"]])
        self.assertIn("| x |  |  |", md)


class OfficeParserTest(unittest.TestCase):
    def test_pptx_preserves_slide_order(self):
        """pptx 语料是操作步骤教程，页码本身就是步骤顺序，丢了就失去上下文。"""
        prs = Presentation()
        for title in ("一、下载并安装APP", "二、打开APP注册账户"):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = title
        buf = BytesIO()
        prs.save(buf)

        result = parse_bytes(".pptx", buf.getvalue())
        self.assertTrue(result.ok)
        self.assertIn("## 第 1 页", result.text)
        self.assertIn("## 第 2 页", result.text)
        self.assertLess(result.text.index("一、下载并安装APP"), result.text.index("二、打开APP注册账户"))

    def test_docx_keeps_paragraph_and_table_in_document_order(self):
        doc = DocxDocument()
        doc.add_paragraph("正文段落一")
        table = doc.add_table(rows=1, cols=2)
        table.cell(0, 0).text = "表头A"
        table.cell(0, 1).text = "表头B"
        doc.add_paragraph("正文段落二")
        buf = BytesIO()
        doc.save(buf)

        result = parse_bytes(".docx", buf.getvalue())
        self.assertTrue(result.ok)
        text = result.text
        self.assertLess(text.index("正文段落一"), text.index("表头A"))
        self.assertLess(text.index("表头A"), text.index("正文段落二"))


@unittest.skipUnless(REAL_DOC.exists(), "语料文件缺失")
class LegacyDocTest(unittest.TestCase):
    """.doc 是 OLE2 二进制格式（已用魔数确认，不是改名的 docx）。"""

    def test_extracts_text_and_flags_degraded(self):
        result = parse_path(REAL_DOC)
        self.assertTrue(result.ok)
        self.assertIn("遗失补办", result.text)
        self.assertTrue(
            result.degraded,
            "旧 .doc 走的是启发式 piece table 解析，必须标记降级并说明，"
            "不能让调用方以为这份文本和其它格式一样可信",
        )
        self.assertTrue(result.reason, "降级时必须给出说明")


@unittest.skipUnless(REAL_XLS.exists(), "语料文件缺失")
class LegacyXlsTest(unittest.TestCase):
    def test_extracts_rows_with_structure(self):
        result = parse_path(REAL_XLS)
        self.assertTrue(result.ok)
        self.assertFalse(result.degraded)
        self.assertIn("大气科学学院", result.text)
        self.assertIn(" | ", result.text, "行内单元格之间要保留分隔，不能塌成一句话")


class OcrAvailabilityTest(unittest.TestCase):
    """OCR 是可选依赖，未安装时必须"显式不可用"，既不崩溃也不静默跳过。

    这正是改造前最糟糕的那种行为：图片被白名单挡掉，使用者完全看不到
    "这里有 4 个流程图没进知识库"。
    """

    @unittest.skipUnless(REAL_JPG.exists(), "语料文件缺失")
    def test_image_result_is_never_silently_empty(self):
        result = parse_path(REAL_JPG)
        # "静默跳过"在这里的具体形态就是：状态说成功、内容却是空的——调用方
        # 无从得知这个文件其实没进知识库。这一条直接禁掉这种组合。
        self.assertFalse(
            result.status is ParseStatus.SUCCESS and not result.text.strip(),
            "不允许出现『SUCCESS 但没有任何文本』这种静默跳过",
        )
        if result.status is ParseStatus.UNAVAILABLE:
            self.assertIn("rapidocr", result.reason.lower())
            self.assertIn("uv sync --extra ocr", result.reason)
        else:
            # 装了 OCR 依赖的环境：应当真的识别出文字
            self.assertTrue(result.ok)
            self.assertTrue(result.text.strip())


if __name__ == "__main__":
    unittest.main()
